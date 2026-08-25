"""Deterministic PowerPoint renderer for the OCL report.

The PPT is a presentation layer over the same AnalysisResult used by the Excel
workbook. It may select, order and format supported evidence, but it must never
recalculate amounts, change materiality, or create unsupported conclusions.
"""
from __future__ import annotations

from decimal import Decimal
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from ocl_agent.schemas import AnalysisResult, AnalysisTable, Finding, ManagementQuestion

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BLUE = RGBColor(0x00, 0x33, 0x8D)
DARK_BLUE = RGBColor(0x00, 0x20, 0x60)
LIGHT_BLUE = RGBColor(0xD9, 0xE2, 0xF3)
GREY = RGBColor(0xE5, 0xE5, 0xE5)
DARK_GREY = RGBColor(0x55, 0x55, 0x55)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xFF, 0xF2, 0xCC)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x61, 0x00)
FONT = "Arial"


def render_report(analysis: AnalysisResult, questions: tuple[ManagementQuestion, ...], output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _remove_default_slides(prs)

    tables = {table.key: table for table in analysis.tables}
    _cover_slide(prs, analysis)
    _key_deal_issues_slide(prs, analysis.findings)
    if annual := tables.get("annual_balance"):
        _annual_snapshot_slide(prs, annual, analysis.findings)
    if movement := tables.get("movement_review"):
        _rollforward_slide(prs, movement)
    if seasonality := tables.get("seasonality"):
        _seasonality_slide(prs, seasonality)
    if monthly := tables.get("monthly_statistics"):
        _monthly_summary_slide(prs, monthly)
    if questions:
        _management_questions_slide(prs, questions, analysis.findings)
    _final_quality_slide(prs, analysis)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _remove_default_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _decorate(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.34), Inches(8.9), Inches(0.45))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = DARK_BLUE
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.56), Inches(0.78), Inches(9.8), Inches(0.35))
        sp = st.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT; sp.font.size = Pt(9); sp.font.color.rgb = DARK_GREY
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.22))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Other Current Liabilities FDD | Source: reconciled OCL databook"
    fp.font.name = FONT; fp.font.size = Pt(7); fp.font.color.rgb = DARK_GREY


def _cover_slide(prs: Presentation, analysis: AnalysisResult) -> None:
    slide = _blank(prs)
    _decorate(slide, "Other Current Liabilities", "Financial due diligence report")
    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(7.2), Inches(1.1))
    p = title.text_frame.paragraphs[0]
    p.text = "Key deal issues in Other Current Liabilities"
    p.font.name = FONT; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK_BLUE
    bullets = [
        f"Annual periods analysed: {', '.join(analysis.annual_periods) or 'not available'}",
        f"Monthly periods analysed: {len(analysis.monthly_periods)}" if analysis.monthly_periods else "Monthly detail not available",
        f"Material findings: {len(analysis.findings)}",
    ]
    _bullet_box(slide, bullets, Inches(0.82), Inches(3.0), Inches(6.1), Inches(1.6), size=13)
    _metric_tiles(slide, [
        ("Findings", len(analysis.findings)),
        ("Questions", "AI-led" if analysis.findings else 0),
        ("Latest FY", analysis.latest_annual_period or "N/A"),
    ], Inches(7.25), Inches(1.55))


def _key_deal_issues_slide(prs: Presentation, findings: tuple[Finding, ...]) -> None:
    slide = _blank(prs)
    _decorate(slide, "Key deal issues", "Commercial headlines before supporting detail")
    top = list(findings[:5])
    if not top:
        _statement(slide, "No material deterministic findings were triggered by the available data.", Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.7))
        return
    rows = []
    for finding in top:
        rows.append([
            finding.priority,
            _short(_deal_issue_title(finding), 38),
            _short(_figure_from_metrics(finding), 25),
            _short(_so_what(finding), 92),
            _short(_ask_management(finding), 82),
        ])
    _table(slide, ["Priority", "Issue", "Figure", "Why it matters", "Question for management"], rows, Inches(0.55), Inches(1.45), Inches(12.2), Inches(4.95), widths=[0.9, 2.25, 1.45, 3.75, 3.85])


def _annual_snapshot_slide(prs: Presentation, table: AnalysisTable, findings: tuple[Finding, ...]) -> None:
    slide = _blank(prs)
    _decorate(slide, "Annual snapshot", "Movement and composition of OCL balances")
    rows = _limit_rows(table.rows, 8)
    _table(slide, table.headers, rows, Inches(0.55), Inches(1.35), Inches(7.6), Inches(4.9))
    bullets = [_short(f.text, 130) for f in findings[:4]] or ["No material annual movement finding was triggered."]
    _side_panel(slide, "Key messages", bullets, Inches(8.45), Inches(1.35), Inches(4.25), Inches(4.9))


def _rollforward_slide(prs: Presentation, table: AnalysisTable) -> None:
    slide = _blank(prs)
    _decorate(slide, "Roll-forward / movement review", "Opening, movements and closing evidence where supported")
    _table(slide, table.headers, _limit_rows(table.rows, 9), Inches(0.55), Inches(1.35), Inches(12.2), Inches(4.9))
    _note(slide, "Movements are shown only when the data-preparation and semantic handoff provide an explicit movement dataset and movement roles.")


def _seasonality_slide(prs: Presentation, table: AnalysisTable) -> None:
    slide = _blank(prs)
    _decorate(slide, "Seasonality", "Year-end representativeness of OCL balances")
    _table(slide, table.headers, _limit_rows(table.rows, 7), Inches(0.55), Inches(1.35), Inches(7.75), Inches(4.8))
    chart_rows = [row for row in table.rows if len(row) >= 4][:6]
    if chart_rows:
        chart_data = CategoryChartData()
        chart_data.categories = [str(row[0]) for row in chart_rows]
        chart_data.add_series("Deviation", [_num(row[3]) for row in chart_rows])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(8.6), Inches(1.65), Inches(3.8), Inches(3.7), chart_data).chart
        chart.has_legend = False
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.category_axis.tick_labels.font.size = Pt(7)
    _note(slide, "Spike / dip flags compare year-end to the trailing 12-month average; unsupported monthly detail does not create this slide.")


def _monthly_summary_slide(prs: Presentation, table: AnalysisTable) -> None:
    slide = _blank(prs)
    _decorate(slide, "Top item monthly summary", "Average, range and latest monthly balance")
    rows = sorted(table.rows, key=lambda r: abs(_num(r[-1])), reverse=True)[:6]
    _table(slide, table.headers, rows, Inches(0.55), Inches(1.35), Inches(7.4), Inches(4.9))
    if rows:
        chart_data = CategoryChartData()
        chart_data.categories = [str(row[0]) for row in rows]
        chart_data.add_series("Latest", [_num(row[-1]) for row in rows])
        chart_data.add_series("Average", [_num(row[1]) for row in rows])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(8.3), Inches(1.55), Inches(4.0), Inches(3.9), chart_data).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(7)
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.category_axis.tick_labels.font.size = Pt(7)


def _management_questions_slide(prs: Presentation, questions: tuple[ManagementQuestion, ...], findings: tuple[Finding, ...]) -> None:
    slide = _blank(prs)
    _decorate(slide, "Questions for management", "Evidence-led questions arising from material findings")
    by_id = {finding.finding_id: finding for finding in findings}
    rows = []
    for idx, item in enumerate(questions[:8], start=1):
        finding = by_id.get(item.linked_finding_id or "")
        rows.append([idx, _theme(finding), _short(item.question, 92), _short(item.rationale, 70), ""])
    _table(slide, ["#", "Theme", "Question", "Evidence / rationale", "Management response"], rows, Inches(0.55), Inches(1.35), Inches(12.2), Inches(5.3), widths=[0.45, 1.7, 4.3, 3.15, 2.6])


def _final_quality_slide(prs: Presentation, analysis: AnalysisResult) -> None:
    slide = _blank(prs)
    _decorate(slide, "Data sources and quality checks", "Traceability and unsupported-analysis discipline")
    bullets = [
        "PPT is generated from the same reconciled AnalysisResult used by the Excel databook.",
        "Slides are created only when the supporting analysis table or finding exists.",
        "Amounts and materiality conclusions are not recalculated in PowerPoint.",
        "Unsupported monthly, roll-forward or seasonality analyses are omitted rather than shown blank.",
    ]
    _bullet_box(slide, bullets, Inches(0.8), Inches(1.55), Inches(11.2), Inches(2.4), size=14)
    _metric_tiles(slide, [("Tables", len(analysis.tables)), ("Findings", len(analysis.findings)), ("Latest FY", analysis.latest_annual_period or "N/A")], Inches(1.0), Inches(4.4))


def _table(slide, headers: Iterable[Any], rows: Iterable[Iterable[Any]], x, y, w, h, widths: list[float] | None = None) -> None:
    headers = [str(header) for header in headers]
    rows = [list(row) for row in rows]
    n_cols = min(len(headers), 7)
    n_rows = max(1, min(len(rows), 10)) + 1
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = shape.table
    if widths:
        for col_idx, width in enumerate(widths[:n_cols]):
            table.columns[col_idx].width = Inches(width)
    for col_idx in range(n_cols):
        cell = table.cell(0, col_idx)
        cell.text = headers[col_idx]
        _cell(cell, fill=BLUE, color=WHITE, bold=True, size=8)
    for row_idx in range(1, n_rows):
        row = rows[row_idx - 1] if row_idx - 1 < len(rows) else []
        for col_idx in range(n_cols):
            value = row[col_idx] if col_idx < len(row) else ""
            cell = table.cell(row_idx, col_idx)
            cell.text = _format_value(value)
            fill = LIGHT_BLUE if row_idx % 2 == 0 else WHITE
            if str(value).upper() in {"HIGH", "FAIL"}:
                fill = AMBER
            _cell(cell, fill=fill, color=BLACK, bold=False, size=7)
            if _is_number(value):
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _cell(cell, *, fill: RGBColor, color: RGBColor, bold: bool, size: int) -> None:
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.alignment = PP_ALIGN.LEFT


def _side_panel(slide, title: str, bullets: list[str], x, y, w, h) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = GREY
    box.line.color.rgb = GREY
    t = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.15), w - Inches(0.35), h - Inches(0.3))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = DARK_BLUE
    for bullet in bullets:
        b = tf.add_paragraph(); b.text = bullet; b.level = 0
        b.font.name = FONT; b.font.size = Pt(9); b.font.color.rgb = BLACK


def _bullet_box(slide, bullets: list[str], x, y, w, h, *, size: int = 12) -> None:
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = BLACK


def _statement(slide, text: str, x, y, w, h) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = GREY
    shape.line.color.rgb = GREY
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT; p.font.size = Pt(14); p.font.color.rgb = DARK_BLUE


def _metric_tiles(slide, metrics: list[tuple[str, Any]], x, y) -> None:
    for idx, (label, value) in enumerate(metrics):
        left = x + Inches(idx * 1.85)
        tile = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, Inches(1.65), Inches(1.0))
        tile.fill.solid(); tile.fill.fore_color.rgb = LIGHT_BLUE
        tile.line.color.rgb = BLUE
        tf = tile.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.name = FONT; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = DARK_BLUE
        s = tf.add_paragraph()
        s.text = label
        s.font.name = FONT; s.font.size = Pt(8); s.font.color.rgb = DARK_GREY


def _note(slide, text: str) -> None:
    shape = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(11.9), Inches(0.35))
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT; p.font.size = Pt(8); p.font.color.rgb = DARK_GREY


def _limit_rows(rows, limit: int):
    return tuple(rows[:limit])


def _format_value(value) -> str:
    if value is None:
        return ""
    if isinstance(value, Decimal):
        return f"{value:,.0f}"
    if isinstance(value, float):
        if abs(value) <= 2:
            return f"{value:.1%}"
        return f"{value:,.1f}"
    if isinstance(value, int):
        return f"{value:,}"
    return str(value)


def _is_number(value) -> bool:
    return isinstance(value, (Decimal, float, int))


def _num(value) -> float:
    if isinstance(value, Decimal):
        return float(value)
    if isinstance(value, (float, int)):
        return float(value)
    try:
        return float(str(value).replace(",", ""))
    except ValueError:
        return 0.0


def _short(text: Any, limit: int) -> str:
    text = str(text or "")
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "…"


def _figure_from_metrics(finding: Finding) -> str:
    metrics = finding.metrics
    for key in ("change", "absolute_value", "latest", "peak_value"):
        if metrics.get(key) is not None:
            try:
                return f"{Decimal(str(metrics[key])):,.0f}"
            except Exception:
                return str(metrics[key])
    if metrics.get("share_pct") is not None:
        return f"{float(metrics['share_pct']):.1f}%"
    return ""


def _deal_issue_title(finding: Finding) -> str:
    return {
        "DEBT_LIKE": "Debt-like items within OCL",
        "DEBT_LIKE_GAP": "FDD vs management debt-like classification gap",
        "ONE_OFF": "One-off / non-recurring OCL items",
        "SEASONALITY": "Year-end balance may not be representative",
        "MONTHLY_VARIABILITY": "Material monthly volatility",
        "STALE_BALANCE": "Potential stale accrual",
        "NEW_ITEM": "New closing obligation",
        "CLIFF": "Balance released or settled to nil",
        "CATEGORY_MOVEMENT": "Category movement requiring explanation",
        "TOTAL_CHANGE": "Total OCL movement",
        "CONCENTRATION": "OCL concentration",
    }.get(finding.finding_type, finding.title)


def _theme(finding: Finding | None) -> str:
    if finding is None:
        return "Other OCL matters"
    return {
        "DEBT_LIKE": "Net debt / equity value",
        "DEBT_LIKE_GAP": "Net debt / equity value",
        "ONE_OFF": "Quality of earnings",
        "SEASONALITY": "Seasonality / phasing",
        "MONTHLY_VARIABILITY": "Seasonality / phasing",
        "STALE_BALANCE": "Working capital validity",
        "NEW_ITEM": "Completeness / validity",
        "CLIFF": "QoE / release risk",
        "CATEGORY_MOVEMENT": "Balance movement",
        "TOTAL_CHANGE": "Balance movement",
        "CONCENTRATION": "Composition",
    }.get(finding.finding_type, "Other OCL matters")


def _so_what(finding: Finding) -> str:
    return {
        "DEBT_LIKE": "Could affect net debt / equity value depending on the transaction definition.",
        "DEBT_LIKE_GAP": "Potential deal-value reconciliation item between management and FDD treatment.",
        "ONE_OFF": "May not represent recurring operating liabilities for normalized working capital.",
        "SEASONALITY": "Year-end working capital may not represent the normal in-year level.",
        "MONTHLY_VARIABILITY": "A single balance-sheet date may not represent the run-rate.",
        "STALE_BALANCE": "May be a valid long-dated obligation or an accrual requiring reassessment.",
        "NEW_ITEM": "New obligation versus prior period; origin and settlement should be evidenced.",
        "CLIFF": "Could reflect settlement, release or reversal affecting QoE / WC interpretation.",
        "CATEGORY_MOVEMENT": "Operational/accounting driver should be understood before conclusions.",
        "TOTAL_CHANGE": "Overall OCL movement should be reconciled to category drivers.",
        "CONCENTRATION": "Largest category drives much of the OCL risk and settlement profile.",
    }.get(finding.finding_type, "Material item requiring evidence-led interpretation.")


def _ask_management(finding: Finding) -> str:
    metrics = finding.metrics
    if finding.finding_type == "SEASONALITY":
        return f"Explain why {metrics.get('category')} differs materially from the 12-month average."
    if finding.finding_type == "CATEGORY_MOVEMENT":
        return f"Explain the principal driver of the movement in {metrics.get('category')}."
    if finding.finding_type == "CONCENTRATION":
        return f"Provide composition and settlement timing for {metrics.get('category')}."
    if finding.finding_type == "TOTAL_CHANGE":
        return "Explain the main operational/accounting drivers of the overall OCL movement."
    return "Provide support for the underlying obligation, classification and expected settlement."
}
