from decimal import Decimal
from pathlib import Path

from pptx import Presentation

from ocl_agent.part4_report.run import run_report
from ocl_agent.schemas import AnalysisResult, AnalysisTable, Finding, ManagementQuestion


def test_ppt_renderer_uses_guided_slide_structure(tmp_path: Path):
    analysis = AnalysisResult(
        findings=(
            Finding(
                "F_TOTAL_CHANGE",
                "OCL movement",
                "Total OCL increased materially from FY24 to FY25.",
                ("FY24", "FY25"),
                "TOTAL_CHANGE",
                {"change": "250000", "previous_period": "FY24", "latest_period": "FY25"},
                "HIGH",
            ),
        ),
        tables=(
            AnalysisTable("annual_balance", "Annual OCL balance by category", ("Category", "FY24", "FY25"), (("Bonus accrual", Decimal("300000"), Decimal("450000")), ("Total OCL", Decimal("600000"), Decimal("850000")))),
            AnalysisTable("movement_review", "Latest annual movement review", ("Category", "FY24", "FY25", "Movement", "Movement %"), (("Bonus accrual", Decimal("300000"), Decimal("450000"), Decimal("150000"), 50.0),)),
            AnalysisTable("seasonality", "Year-end representativeness / seasonality", ("Category", "12M Average", "Year End", "Deviation", "Deviation %", "Flag"), (("Bonus accrual", Decimal("350000"), Decimal("450000"), Decimal("100000"), 28.6, "YEAR-END SPIKE"),)),
            AnalysisTable("monthly_statistics", "Monthly OCL statistics by category", ("Category", "Average", "Minimum", "Maximum", "Std_Dev", "Latest"), (("Bonus accrual", Decimal("350000"), Decimal("250000"), Decimal("450000"), Decimal("55000"), Decimal("450000")),)),
        ),
        annual_periods=("FY24", "FY25"),
        monthly_periods=tuple(f"2025-{month:02d}" for month in range(1, 13)),
        latest_annual_period="FY25",
    )
    questions = (
        ManagementQuestion("Q1", "Please explain the main driver of the OCL movement.", "The annual movement was material.", ("F_TOTAL_CHANGE",), "F_TOTAL_CHANGE", "HIGH"),
    )

    report = run_report(analysis, questions, tmp_path)
    prs = Presentation(report)
    assert len(prs.slides) == 8
    assert prs.slide_width == 12191875
    assert prs.slide_height == 6858000

    slide_text = ["\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")) for slide in prs.slides]
    assert "Key deal issues in Other Current Liabilities" in slide_text[0]
    assert "Key deal issues" in slide_text[1]
    assert "Annual snapshot" in slide_text[2]
    assert "Roll-forward / movement review" in slide_text[3]
    assert "Seasonality" in slide_text[4]
    assert "Questions for management" in slide_text[6]

    for slide in prs.slides:
        assert any("Other Current Liabilities FDD" in getattr(shape, "text", "") for shape in slide.shapes)
